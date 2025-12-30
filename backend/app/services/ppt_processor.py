"""
PPT 处理服务
负责将 PPT/PDF 文件转换为图片，并提取文本内容
"""

import os
import tempfile
from typing import List, Tuple
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
import PyPDF2


class PPTProcessor:
    """PPT/PDF 处理器"""

    def __init__(self, output_dir: str):
        """
        初始化处理器

        Args:
            output_dir: 输出目录，用于保存转换后的图片
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def process_file(self, file_path: str, file_type: str) -> List[Tuple[str, str]]:
        """
        处理上传的文件

        Args:
            file_path: 文件路径
            file_type: 文件类型 (pdf, ppt, pptx)

        Returns:
            List of (image_path, text_content) tuples
        """
        if file_type == 'pdf':
            return self._process_pdf(file_path)
        elif file_type in ['ppt', 'pptx']:
            return self._process_pptx(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {file_type}")

    def _process_pdf(self, pdf_path: str) -> List[Tuple[str, str]]:
        """
        处理 PDF 文件

        Args:
            pdf_path: PDF 文件路径

        Returns:
            List of (image_path, text_content) tuples
        """
        results = []

        # 转换 PDF 为图片
        try:
            images = convert_from_path(pdf_path, dpi=200)
        except Exception as e:
            raise RuntimeError(f"PDF 转图片失败: {str(e)}")

        # 提取 PDF 文本
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                texts = []
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    texts.append(text if text else "")
        except Exception as e:
            print(f"PDF 文本提取失败: {str(e)}")
            texts = [""] * len(images)

        # 保存图片
        for idx, image in enumerate(images):
            slide_number = idx + 1
            image_filename = f"slide_{slide_number}.png"
            image_path = self.output_dir / image_filename

            # 保存图片
            image.save(str(image_path), 'PNG')

            # 获取对应的文本
            text_content = texts[idx] if idx < len(texts) else ""

            results.append((str(image_path), text_content))

        return results

    def _process_pptx(self, pptx_path: str) -> List[Tuple[str, str]]:
        """
        处理 PPTX 文件

        Args:
            pptx_path: PPTX 文件路径

        Returns:
            List of (image_path, text_content) tuples

        Note:
            对于 .ppt 文件，需要先转换为 .pptx
            或者使用 LibreOffice 转换为 PDF 后再处理
        """
        # 如果是 .ppt 文件，先转换为 PDF
        if pptx_path.endswith('.ppt'):
            return self._convert_ppt_to_pdf_and_process(pptx_path)

        # 处理 .pptx 文件
        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            raise RuntimeError(f"无法打开 PPTX 文件: {str(e)}")

        results = []

        # 提取文本
        texts = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            texts.append("\n".join(slide_text))

        # 对于 PPTX，我们需要将其转换为 PDF，然后再转为图片
        # 这需要 LibreOffice 或其他工具
        # 尝试使用 LibreOffice 转换
        pdf_path = self._convert_pptx_to_pdf(pptx_path)

        if pdf_path:
            # LibreOffice 转换成功，使用 PDF 转图片
            images = convert_from_path(pdf_path, dpi=200)

            for idx, image in enumerate(images):
                slide_number = idx + 1
                image_filename = f"slide_{slide_number}.png"
                image_path = self.output_dir / image_filename

                # 保存图片
                image.save(str(image_path), 'PNG')

                # 获取对应的文本
                text_content = texts[idx] if idx < len(texts) else ""

                results.append((str(image_path), text_content))

            # 清理临时 PDF
            os.remove(pdf_path)
        else:
            # LibreOffice 不可用，使用备用方案：生成简单的文本图片
            from PIL import Image, ImageDraw, ImageFont
            import textwrap

            for idx, slide_text in enumerate(texts):
                slide_number = idx + 1
                image_filename = f"slide_{slide_number}.png"
                image_path = self.output_dir / image_filename

                # 创建一个 1920x1080 的白色背景图片
                img = Image.new('RGB', (1920, 1080), color='white')
                draw = ImageDraw.Draw(img)

                # 尝试使用系统字体，如果失败则使用默认字体
                try:
                    # Windows 系统字体路径
                    font_large = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 48)
                    font_small = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 32)
                except:
                    font_large = ImageFont.load_default()
                    font_small = ImageFont.load_default()

                # 绘制标题
                title = f"Slide {slide_number}"
                draw.text((100, 100), title, fill='black', font=font_large)

                # 绘制文本内容（换行处理）
                if slide_text:
                    wrapped_text = textwrap.fill(slide_text, width=60)
                    draw.text((100, 200), wrapped_text, fill='#333333', font=font_small)
                else:
                    draw.text((100, 200), "(No text content)", fill='#999999', font=font_small)

                # 保存图片
                img.save(str(image_path), 'PNG')

                results.append((str(image_path), slide_text))

            import logging
            logger = logging.getLogger(__name__)
            logger.warning("Using fallback text-to-image conversion (LibreOffice not available)")

        return results

    def _convert_pptx_to_pdf(self, pptx_path: str) -> str:
        """
        将 PPTX 转换为 PDF（使用 LibreOffice）

        Args:
            pptx_path: PPTX 文件路径

        Returns:
            PDF 文件路径
        """
        try:
            import subprocess
            import sys

            # 创建临时目录
            temp_dir = tempfile.mkdtemp()
            output_pdf = os.path.join(temp_dir, "output.pdf")

            # 在 Windows 上设置正确的编码环境
            env = os.environ.copy()
            if sys.platform == 'win32':
                env['PYTHONIOENCODING'] = 'utf-8'

            # 使用 LibreOffice 转换
            # 注意：需要系统安装 LibreOffice
            result = subprocess.run([
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                pptx_path
            ], check=True, capture_output=True, text=True, encoding='utf-8', errors='replace', env=env)

            # 查找生成的 PDF 文件
            pdf_files = list(Path(temp_dir).glob('*.pdf'))
            if pdf_files:
                return str(pdf_files[0])
            else:
                raise RuntimeError("Converted PDF file not found")

        except FileNotFoundError as e:
            # LibreOffice not installed
            import logging
            logger = logging.getLogger(__name__)
            logger.error("LibreOffice not installed, cannot convert PPTX to PDF")
            return None
        except subprocess.CalledProcessError as e:
            import logging
            logger = logging.getLogger(__name__)
            # 避免直接打印可能包含非 ASCII 字符的错误信息
            logger.error(f"PPTX to PDF conversion failed with exit code: {e.returncode}")
            return None
        except Exception as e:
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"PPTX to PDF conversion error: {type(e).__name__}")
            return None

    def _convert_ppt_to_pdf_and_process(self, ppt_path: str) -> List[Tuple[str, str]]:
        """
        Convert .ppt file to PDF and process

        Args:
            ppt_path: .ppt file path

        Returns:
            List of (image_path, text_content) tuples
        """
        pdf_path = self._convert_pptx_to_pdf(ppt_path)
        if pdf_path:
            results = self._process_pdf(pdf_path)
            os.remove(pdf_path)
            return results
        else:
            raise RuntimeError("Cannot convert .ppt file, please convert to .pptx or PDF first")


def get_file_type(filename: str) -> str:
    """
    Determine file type from filename

    Args:
        filename: filename

    Returns:
        file type (pdf, ppt, pptx)
    """
    ext = filename.lower().split('.')[-1]
    if ext in ['pdf', 'ppt', 'pptx']:
        return ext
    else:
        raise ValueError(f"Unsupported file type: {ext}")
